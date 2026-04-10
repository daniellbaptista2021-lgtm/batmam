"""Gerador de apresentacoes PPTX."""
from __future__ import annotations
import json
from pathlib import Path
from .base import STATIC_DIR, ask_ai, unique_name, file_url

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def generate(prompt: str) -> dict:
    system = """Voce e um designer de apresentacoes. Gere slides em JSON.
Formato EXATO (sem markdown, sem explicacao):
{
  "title": "Titulo da Apresentacao",
  "slides": [
    {
      "title": "Slide 1 - Titulo",
      "content": ["Topico 1", "Topico 2", "Topico 3"],
      "notes": "Notas do apresentador"
    }
  ]
}
- Minimo 6 slides
- Primeiro slide: titulo e subtitulo
- Ultimo slide: agradecimento/contato
- Conteudo profissional em portugues brasileiro
- Retorne APENAS o JSON valido"""

    model = "deepseek-chat"
    raw = ask_ai(prompt, system=system, model=model, max_tokens=4096)

    text = raw.strip()
    if text.startswith("```"):
        text = "\n".join(text.split("\n")[1:])
    if text.endswith("```"):
        text = text[:-3]
    text = text.strip()

    data = json.loads(text)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    PURPLE = RGBColor(0x7C, 0x3A, 0xED)
    DARK = RGBColor(0x1A, 0x1A, 0x25)
    WHITE = RGBColor(0xF0, 0xF0, 0xF5)
    GRAY = RGBColor(0x9D, 0x9D, 0xB5)

    slides = data.get("slides", [])

    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = DARK

        title_text = slide_data.get("title", "")
        content_items = slide_data.get("content", [])
        notes_text = slide_data.get("notes", "")

        # Title
        left = Inches(0.8)
        top = Inches(0.6) if i == 0 else Inches(0.5)
        width = Inches(11.7)
        height = Inches(1.5) if i == 0 else Inches(1)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36) if i == 0 else Pt(28)
        p.font.bold = True
        p.font.color.rgb = PURPLE if i == 0 else WHITE
        p.alignment = PP_ALIGN.LEFT if i > 0 else PP_ALIGN.CENTER

        # Content bullets
        if content_items and i > 0:
            c_top = Inches(1.8)
            c_height = Inches(5)
            txBox2 = slide.shapes.add_textbox(left, c_top, width, c_height)
            tf2 = txBox2.text_frame
            tf2.word_wrap = True

            for j, item in enumerate(content_items):
                p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                p.text = f"  {item}"
                p.font.size = Pt(18)
                p.font.color.rgb = GRAY
                p.space_after = Pt(12)

        # Accent line
        if i > 0:
            from pptx.shapes.autoshape import Shape
            shape = slide.shapes.add_shape(
                1, left, Inches(1.5), Inches(2), Inches(0.04)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = PURPLE
            shape.line.fill.background()

        # Notes
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    filename = unique_name(data.get("title", prompt[:30]), ".pptx")
    out_dir = STATIC_DIR / "files"
    out_dir.mkdir(parents=True, exist_ok=True)
    filepath = out_dir / filename
    prs.save(filepath)

    url = file_url(f"static/files/{filename}")
    size = filepath.stat().st_size

    return {
        "type": "pptx",
        "name": filename,
        "url": url,
        "size": size,
        "path": str(filepath),
    }
